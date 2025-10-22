import os 
from pptx import Presentation
from pptx.util import Pt
import re,math
from PIL import ImageFont
from pptx.dml.color import RGBColor

font_bold = ImageFont.truetype("./fonts/poppins/Poppins-Bold.ttf", size=80)

SECTIONS = ["HOOK", "HOOK_SUB", "STORY", "INSIGHT", "VALUE", "CTA"]

def parse_post(post_text: str):
    """Parse post sections like [HOOK], [STORY], etc."""
    parts = {}
    current = None
    for line in post_text.splitlines():
        line = line.strip()
        tag = line.strip("[]").upper()
        if tag in SECTIONS:
            current = tag
            parts[current] = []
        elif current and line:
            parts[current].append(line)
    return parts

def hex_to_rgb(hex_color:str = "#000"):
    """
    Converte cor hexadecimal para RGBColor.
    Aceita: "#fff", "#ffffff", "fff", "ffffff"
    """
    hex_color = hex_color.lstrip('#')
    
    # Se for formato curto (#fff), expande para #ffffff
    if len(hex_color) == 3:
        hex_color = ''.join([c*2 for c in hex_color])
    
    # Converte para RGB
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    
    return RGBColor(r, g, b)

def dynamic_font_size_simple(text, max_width_pt=505, max_height_pt=395,
                              base_line_spacing=0.85, min_font_size=17, 
                              max_font_size=110, max_caracteres=840):
    """
    Versão simplificada que usa estimativas matemáticas.
    Mais rápida mas menos precisa.
    MAX x caracteres, based on the box size
    Elegant logarithmic font scaling tuned for visual balance:
      • 25 chars → ~110 pt
      • 50 chars → ~82 pt
      • 80 chars → ~50 pt
      • 110 chars → ~42 pt
    """
    num_chars = len(text.strip())
    print(num_chars)
    words = text.split()
    avg_word_length = num_chars / max(1, len(words))
    
    # Estimativa de caracteres por linha (0.5 é a proporção média char_width/font_size)
    def chars_per_line(font_size):
        avg_char_width = font_size * (1.1 if num_chars > 120 else 1.25) # adjust to big letters
        print(avg_char_width)
        return int(max_width_pt / avg_char_width)
    
    # Busca binária
    low, high = min_font_size, max_font_size
    best_size = min_font_size
    
    for _ in range(20):
        mid = (low + high) / 2
        cpl = chars_per_line(mid)
        estimated_lines = math.floor(num_chars / cpl)
        estimated_height = estimated_lines * mid * base_line_spacing
        if estimated_height <= max_height_pt:
            best_size = mid
            low = mid + 1
            print(high)
        else:
            high = mid - 1
    
    # Ajuste fino do line spacing
    line_spacing = base_line_spacing
    print(base_line_spacing)
    if num_chars > 100:
        line_spacing = min(1.0, base_line_spacing + 0.1)
    # else: 
    #     line_spacing = min(1.0, base_line_spacing + 0.05)
    
    return math.floor(best_size), round(line_spacing, 2), 1 if num_chars > max_caracteres else 0

def dynamic_font_size_and_spacing(text, base_line_height=0.8):
    """
    Elegant logarithmic font scaling tuned for visual balance:
      • 25 chars → ~110 pt
      • 50 chars → ~82 pt
      • 80 chars → ~50 pt
      • 110 chars → ~42 pt
    """
    num_chars = max(1, len(text.strip()))
    print(len(text))
    print(num_chars)

    # Tuned constants (steeper decay)
    A = 160   # intercept — base size for very short text
    B = 25    # decay rate constant

    # Smooth logarithmic decay
    font_size_pt = A - B * math.log(num_chars-10)

    # Clamp for safety
    font_size_pt = max(25, min(font_size_pt, 110))

    # Line height increases gently for long text
    line_spacing = min(0.9, base_line_height + (num_chars / 60) * 0.15)

    return round(font_size_pt, 2), round(line_spacing, 2)

def fill_carousel(post_text, template_path="brand_carousel_template_portrait.pptx", output_path="carousel_filled.pptx"):
    prs = Presentation(template_path)
    parts = parse_post(post_text)
    slides = prs.slides

    # 🎯 Font and style settings
    FONT_SIZES = {
        "HOOK": Pt(85),
        "HOOK_SUB": Pt(17),
        "STORY": Pt(48),
        "INSIGHT": Pt(28),
        "VALUE": Pt(28),
        "CTA": Pt(36),
    }

    LINE_HEIGHTS = {
        "HOOK": 0.7,      # reduz espaçamento entre linhas
        "HOOK_SUB": 1.0,
        "STORY": 1.2,
        "INSIGHT": 1.15,
        "VALUE": 1.1,
        "CTA": 1.1,
    }
 # -------------------------------------
    # SLIDE 0 — COVER (HOOK + HOOK_SUB)
    # -------------------------------------
    if "HOOK" in parts or "HOOK_SUB" in parts:
        new_hook = "\n".join(parts.get("HOOK", []))
        new_hook_sub = "\n".join(parts.get("HOOK_SUB", []))
        hook_text = new_hook.strip()

        # 🧠 Dynamically calculate font size and line spacing
        font_size, line_spacing, passed = dynamic_font_size_simple(
                hook_text,
                max_width_pt=395,
                max_height_pt=480,
                max_caracteres=840
            )
        print("FALSE" if passed else "True")

        for shape in slides[0].shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text

            # 🔍 Replace placeholders within the same shape
            new_text = text
            if "[HOOK]" in new_text:
                new_text = new_text.replace("[HOOK]", new_hook)
            if "[HOOK_SUB]" in new_text:
                new_text = new_text.replace("[HOOK_SUB]", new_hook_sub)

            # 🖋️ If replacements occurred → apply styles
            if new_text != text:
                shape.text_frame.text = new_text

                for paragraph in shape.text_frame.paragraphs:
                    # Apply dynamic line height
                    paragraph.line_spacing = line_spacing

                    for run in paragraph.runs:
                        # HOOK (bold + dynamic font)
                        if new_hook in run.text:
                            run.font.size = Pt(font_size)
                            run.font.bold = True
                            run.font.name = "Poppins"
                            # run.font.color.rgb = hex_to_rgb("#fff")
                        # HOOK_SUB (smaller + normal weight)
                        elif new_hook_sub in run.text:
                            # 70% of hook size, minimum 18pt
                            sub_size = max(18, 20)
                            print(sub_size)
                            run.font.size = Pt(sub_size)
                            run.font.bold = False
                            run.font.name = "Poppins thin"
                 
                            # Subtitles get slightly looser line height
                            paragraph.line_spacing = min(1.2, line_spacing + 0.15)

    # # -------------------------------------
    # # SLIDE 1 — STORY
    # # -------------------------------------
    # if "STORY" in parts:
    #     story_text = "\n".join(parts["STORY"])
    #     for shape in slides[1].shapes:
    #         if shape.has_text_frame:
    #             shape.text_frame.text = story_text
    #             for p in shape.text_frame.paragraphs:
    #                 p.line_spacing = LINE_HEIGHTS["STORY"]
    #                 for run in p.runs:
    #                     run.font.size = FONT_SIZES["STORY"]

    # # -------------------------------------
    # # SLIDE 2 — INSIGHT + VALUE
    # # -------------------------------------
    # if "INSIGHT" in parts:
    #     slides[2].shapes[0].text = "Key Insights"
    #     for shape in slides[2].shapes:
    #         if shape.has_text_frame and "[INSIGHT]" in shape.text_frame.text:
    #             shape.text_frame.text = "\n".join(parts["INSIGHT"])
    #             for p in shape.text_frame.paragraphs:
    #                 p.line_spacing = LINE_HEIGHTS["INSIGHT"]
    #                 for run in p.runs:
    #                     run.font.size = FONT_SIZES["INSIGHT"]

    # if "VALUE" in parts:
    #     bullet_slide = slides[2]
    #     for shape in bullet_slide.shapes:
    #         if shape.has_text_frame and "✔" in shape.text_frame.text:
    #             shape.text_frame.clear()
    #             for line in parts["VALUE"]:
    #                 p = shape.text_frame.add_paragraph()
    #                 p.text = f"✔ {line}"
    #                 p.line_spacing = LINE_HEIGHTS["VALUE"]
    #                 for run in p.runs:
    #                     run.font.size = FONT_SIZES["VALUE"]

    # # -------------------------------------
    # # SLIDE 3 — CTA
    # # -------------------------------------
    # if "CTA" in parts:
    #     for shape in slides[3].shapes:
    #         if shape.has_text_frame and "[CTA]" in shape.text_frame.text:
    #             shape.text_frame.text = "\n".join(parts["CTA"])
    #             for p in shape.text_frame.paragraphs:
    #                 p.line_spacing = LINE_HEIGHTS["CTA"]
    #                 for run in p.runs:
    #                     run.font.size = FONT_SIZES["CTA"]

    prs.save(output_path)
    print(f"✅ Carousel generated: {output_path}")

# Discover how data-driven B2B marketing can elevate your strategy, delivering insights that shape impactful decisions.

if __name__ == "__main__":
    example_post = """
[HOOK]
Optimize resources, maximize ROI, and extend your B2B marketing budget through data-backed decisions.Optimize resources, maximize ROI, and extend your B2B marketing budget through data-backed decisions.Optimize resources, maximize ROI, and extend your B2B marketing budget through data-backed decisions.Optimize resources, maximize ROI, and extend your B2B marketing budget through data-backed decisions.Optimize resources, maximize ROI, and extend your B2B marketing budget through data-backed decisions.Optimize resources, maximize ROI, and extend your B2B marketing budget through data-backed decisions.
[HOOK_SUB]
Efficiency > Size. For production, cost and speed matter more.
[STORY]
I tested small LLMs on my laptop. They surprised me by matching larger models in several tasks.

[INSIGHT]
Efficiency > Size. For production, cost and speed matter more.

[VALUE]
Define the task clearly
Benchmark small models first
Scale only if necessary

[CTA]
Would you deploy a small model in production?
"""
    fill_carousel(
        example_post,
        template_path="brand_hook.pptx",
        output_path="example_carousel_filled.pptx"
    )



# 🟢 Short hooks (25–40 characters) — should appear large (~90–110 pt)

# AI that fits on your laptop

# Smarter code starts with context

# Where data meets creativity

# Engineering ideas that scale

# From models to meaning

# 🟡 Medium hooks (45–60 characters) — balanced size (~70–85 pt)

# The Power of Data-Driven Marketing in B2B

# How small language models are changing AI forever


# Redefining innovation through efficiency and trust

# Turning everyday data into strategic advantage

# Bringing human insight back to machine learning

# 🟠 Long hooks (70–90 characters) — smaller text (~55–65 pt)

# Discover how local AI models can make enterprise innovation faster and cheaper

# Rethinking leadership in the age of automation and intelligent decision-making

# Why personalization will define the next generation of customer experiences

# Exploring how small teams can compete with big tech using smart AI tools

# Designing systems that connect performance, empathy, and trust

# 🔵 Very long hooks (100–120 characters) — compact (~40–50 pt)

# Discover how data-driven B2B marketing can elevate your strategy, delivering insights that shape impactful decisions.

# Building scalable AI platforms that connect data, design, and decision-making for smarter digital transformation.

# Empowering teams to innovate responsibly by combining human judgment with machine intelligence.

# Transform your workflow with lightweight AI models that bring enterprise-level performance to local environments.

# Exploring how transparency, trust, and accessibility will redefine success in the AI era.


# Perfect 👌 — you want to refine the curve so that:

# Range	Characters	Target Font Size
# 🟢 Short	~25	~105–110 pt
# 🟡 Medium	~50	~75–85 pt
# 🟠 Long	~70–90	~40–55 pt ✅
# 🔵 Very Long	~100–120	~38–45 pt


