import os
import sys
import shutil
from pptx import Presentation

def split_pptx_by_layout(input_dir: str):
    """
    Split each PPTX into separate files by slide layout.
    Each output file contains the slide exactly as it appeared in the original.
    """

    if not os.path.exists(input_dir):
        raise FileNotFoundError(f"Directory not found: {input_dir}")

    for filename in os.listdir(input_dir):
        if not filename.endswith(".pptx"):
            continue
        if "-" not in filename:
            print(f"⚠️ Skipped (invalid filename): {filename}")
            continue

        style, _ = filename.replace(".pptx", "").split("-", 1)
        filepath = os.path.join(input_dir, filename)

        style_dir = os.path.join(input_dir, style)
        os.makedirs(style_dir, exist_ok=True)

        print(f"🔎 Processing {filename} (style={style})")

        prs = Presentation(filepath)

        for idx, slide in enumerate(prs.slides):
            layout_name = slide.slide_layout.name.strip().replace(" ", "_")
            # Export strategy: copy the entire .pptx file and then delete other slides
            tmp_copy = os.path.join(style_dir, f"__tmp_{layout_name}_{idx}.pptx")
            shutil.copy(filepath, tmp_copy)

            tmp_prs = Presentation(tmp_copy)

            # Delete all slides except the one we want
            for i in reversed(range(len(tmp_prs.slides))):
                if i != idx:
                    r_id = tmp_prs.slides._sldIdLst[i].rId
                    tmp_prs.part.drop_rel(r_id)
                    del tmp_prs.slides._sldIdLst[i]

            out_path = os.path.join(style_dir, f"{layout_name}.pptx")
            tmp_prs.save(out_path)
            os.remove(tmp_copy)

            print(f"   ✅ Slide {idx+1} → {out_path}")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python split_templates.py ./templates/blue-blur")
        sys.exit(1)

    input_dir = sys.argv[1]
    split_pptx_by_layout(input_dir)