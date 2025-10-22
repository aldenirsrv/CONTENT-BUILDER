from pptx import Presentation
from pptx.util import Pt
import os
import time
import copy, math, re
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from merge_templates import merge_pptx_slides
from pptx.parts.image import ImagePart
import datetime

"""
LinkedIn Carousel Generator
---------------------------

This script automates the generation of LinkedIn carousel slides by filling PowerPoint 
templates with structured content extracted from a LinkedIn-style post. 
It supports flexible mappings of placeholders (e.g. [HOOK], [STORY], [CTA]) 
to user-provided content blocks, while applying per-block style overrides 
(font size, bold, color, alignment, line spacing).

Author: <your name>
"""

from pptx import Presentation
from pptx.util import Pt
import copy, math, re
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Supported post sections. Any block outside these tags will be ignored.
SECTIONS = {
    "HOOK", 
    "HOOK_SUB", 
    "STORY", 
    "STORY_SUB",
    "TOPIC", 
    "TOPIC_SUB", 
    "CTA", 
    "SUBJECT", 
    "ONE",
    "ONE_SUB",
    "IMAGE_TOP",
    "IMAGE_TOP_SUB",
    "IMAGE_BOTTOM",
    "IMAGE_BOTTOM_SUB",
    "IMAGE_BOTTOM_RIGHT",
    "IMAGE_BOTTOM_RIGHT_SUB",
    "IMAGE_BOTTOM_RIGHT_CAP",
    "IMAGE_BOTTOM_LEFT",
    "IMAGE_BOTTOM_LEFT_SUB",
    "IMAGE_BOTTOM_LEFT_CAP",
    "CTA",
    "CTA_SUB"
}


# Contador global para garantir unicidade
_image_counter = 0

def parse_post(post_text: str):
    """
    Parse LinkedIn-like post text into structured parts.

    Args:
        post_text (str): The full post text containing section tags like [HOOK], [STORY].

    Returns:
        dict: Mapping from section name (e.g., 'HOOK') to list of lines.
    """
    parts = {}
    current = None
    for line in post_text.splitlines():
        line = line.strip()
        if not line:
            continue

        if line.startswith("[") and line.endswith("]"):
            tag = line.strip("[]").upper()
            if tag in SECTIONS:
                current = tag
                parts[current] = []
        elif current:
            parts[current].append(line)
    return parts


def dynamic_font_size_simple(
    text,
    max_width_pt=505,
    max_height_pt=200,
    base_line_spacing=0.85,
    min_font_size=17,
    max_font_size=110,
    max_caracteres=840
):
    """
    Estimate an appropriate font size and line spacing based on content length 
    and bounding box constraints. Uses binary search for efficiency.

    Args:
        text (str): Content to measure.
        max_width_pt (int): Width of the bounding box in points.
        max_height_pt (int): Height of the bounding box in points.
        base_line_spacing (float): Base line spacing multiplier.
        min_font_size (int): Minimum font size allowed.
        max_font_size (int): Maximum font size allowed.
        max_caracteres (int): Hard limit of characters before overflow.

    Returns:
        tuple: (font_size: int, line_spacing: float, overflow_flag: int)
    """
    num_chars = len(text.strip())
    words = text.split()
    avg_word_length = num_chars / max(1, len(words))

    def chars_per_line(font_size):
        # Approximate average char width
        avg_char_width = font_size * (1.1 if num_chars > 120 else 1.25)
        return int(max_width_pt / avg_char_width)

    # Binary search between min and max font size
    low, high = min_font_size, max_font_size
    best_size = min_font_size
    for _ in range(20):
        mid = (low + high) / 2
        cpl = chars_per_line(mid)
        estimated_lines = math.floor(num_chars / cpl)
        estimated_height = estimated_lines * mid * base_line_spacing
        if estimated_height <= max_height_pt:
            best_size = mid
            low = mid + 1
        else:
            high = mid - 1

    # Adjust line spacing slightly for long content
    line_spacing = base_line_spacing
    if num_chars > 100:
        line_spacing = min(1.0, base_line_spacing + 0.1)

    return math.floor(best_size), round(line_spacing, 2), 1 if num_chars > max_caracteres else 0


def hex_to_rgb(hex_color: str = "#000"):
    """
    Convert hex string to RGBColor. Accepts short (#fff) or long (#ffffff) formats.

    Args:
        hex_color (str): Hex color string.

    Returns:
        RGBColor: pptx-compatible RGB color.
    """
    hex_color = hex_color.lstrip('#')
    if len(hex_color) == 3:  # short format (#fff)
        hex_color = ''.join([c*2 for c in hex_color])

    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)
def apply_text_to_slide(prs, placeholder_map, text_parts, output_path, img_path:str = None):
    """
    Apply parsed content into a presentation template.
    """
    global _image_counter
    
    default_img = img_path
    norm_map = {}
    for ph, cfg in placeholder_map.items():
        if isinstance(cfg, str):
            norm_map[ph] = {"key": cfg}
        else:
            norm_map[ph] = dict(cfg)

    ph_pattern = re.compile("(" + "|".join(map(re.escape, norm_map.keys())) + ")")
    
    for slide in prs.slides:
        for shape in slide.shapes:
            # 🎯 Caso 1: Placeholder de imagem
            if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                sp = shape.element
                parent = sp.getparent()
                idx = parent.index(sp)

                parent.remove(sp)

                if os.path.exists(default_img):
                    # Adiciona a nova imagem
                    new_shape = slide.shapes.add_picture(default_img, left, top, width, height)
                    
                    # MODIFICAÇÃO DIRETA: Renomeia a image part
                    _image_counter += 1
                    timestamp = int(time.time() * 1000000)
                    unique_suffix = f"{timestamp}_{_image_counter}"
                    ext = os.path.splitext(default_img)[1]
                    
                    # Acessa a image part e modifica o nome
                    try:
                        # Pega a última relação (a imagem que acabamos de adicionar)
                        rels = slide.part.rels
                        last_rel_id = max([int(rid.replace('rId', '')) for rid in rels.keys()])
                        last_rel = rels[f'rId{last_rel_id}']
                        image_part = last_rel.target_part
                        
                        # Modifica o partname
                        from pptx.opc.packuri import PackURI
                        new_partname = PackURI(f"/ppt/media/image_{unique_suffix}{ext}")
                        image_part._partname = new_partname
                        
                        print(f"🖼️ Imagem criada: image_{unique_suffix}{ext}")
                    except Exception as e:
                        print(f"⚠️ Erro ao renomear: {e}")
                    
                    # Move a imagem para a posição original
                    new_sp = new_shape.element
                    parent.remove(new_sp)
                    parent.insert(idx, new_sp)
                    
                    time.sleep(0.001)  # Pequeno delay

                else:
                    print(f"⚠️ Imagem padrão não encontrada em {default_img}")
                continue

            # Resto do código continua igual...
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                print("🖼️ Shape já é uma imagem")
                continue

            if not shape.has_text_frame:
                continue

            text = shape.text_frame.text
            if not any(ph in text for ph in norm_map):
                continue
            
            tokens = ph_pattern.split(text)
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]

            for tok in tokens:
                if not tok:
                    continue

                if tok in norm_map:
                    cfg = norm_map[tok]
                    key = cfg.get("key")
                    new_text = "\n".join(text_parts.get(key, []))
                    if not new_text:
                        continue

                    font_size, auto_line_spacing, _ = dynamic_font_size_simple(
                        new_text,
                        max_width_pt=395,
                        max_height_pt=cfg.get("text-block-height", 505),
                        max_caracteres=840
                    )
                    
                    size_to_use = cfg.get("size", font_size)
                    spacing_to_use = cfg.get("line_spacing", auto_line_spacing)

                    run = p.add_run()
                    run.text = new_text
                    run.font.size = Pt(max(12, size_to_use))
                    run.font.name = "Poppins" if cfg.get("bold", False) else "Poppins thin"

                    run.font.bold = bool(cfg.get("bold", False))
                    if "font" in cfg:
                        run.font.name = cfg["font"]
                    if "color" in cfg:
                        run.font.color.rgb = hex_to_rgb(cfg["color"])
                    
                    p.line_spacing = cfg.get("line-height", min(1.2, spacing_to_use + 0.15))
                    align_val = cfg.get("align")
                    if align_val:
                        if align_val.lower() == "left":
                            p.alignment = PP_ALIGN.LEFT
                        elif align_val.lower() == "center":
                            p.alignment = PP_ALIGN.CENTER
                        elif align_val.lower() == "right":
                            p.alignment = PP_ALIGN.RIGHT

                else:
                    run = p.add_run()
                    run.text = tok

    prs.save(f"./concluded/{output_path}")
    print(f"💾 Saved: {output_path}")
def build_carousel(post_text, template_mappings, output_path):
    """
    Build carousels from multiple templates. Each template generates a separate PPTX.

    Args:
        post_text (str): Full post text with section tags.
        template_mappings (list): List of dicts with keys:
            - 'template': path to template file
            - 'blocks': mapping placeholders → config
        output_path (str): Base output filename (prefix will be added per template).
    """
    parts = parse_post(post_text)

    for idx, mapping in enumerate(template_mappings, start=1):
        prs = Presentation(mapping["template"])
        file_name = f"{idx}-{output_path}"
        apply_text_to_slide(prs, mapping["blocks"], parts, file_name, mapping["image"])

template_mappings = [
    {"template": "./templates/blue-blur/dark/Cover.pptx",
    "image": None,
     "blocks": {
         "[HOOK]": {"key": "HOOK", "bold": True, "text-block-height": 430},
         "[HOOK_SUB]": {"key": "HOOK_SUB", "size": 17, "bold": False}
     }},

    {"template": "./templates/blue-blur/dark/Story.pptx",
      "image": None,
     "blocks": {
        "[STORY]": {"key": "STORY", "bold": True, "align": "center", "text-block-height": 300},
        "[STORY_SUB]": {"key": "STORY_SUB", "size": 17, "bold": False},
         "[SUBJECT]": {"key": "SUBJECT", "size": 12, "bold": False}
     }},

    {"template": "./templates/blue-blur/dark/Topic.pptx",
      "image": None,
     "blocks": {
        "[TOPIC]": {"key": "TOPIC", "bold": True, "size": 40, "line-height": 0.9},
        "[TOPIC_SUB]": {"key": "TOPIC_SUB", "size": 20, "bold": False, "line-height": 1.5},
        "[SUBJECT]": {"key": "SUBJECT", "size": 12, "bold": False}
     }},

    {
    "template": "./templates/blue-blur/dark/image-top.pptx",
    "image": "./templates/image-w.png",
     "blocks": {
        "[IMAGE_TOP]": {"key": "IMAGE_TOP", "bold": True, "size": 40, "line-height": 0.9},
        "[IMAGE_TOP_SUB]": {"key": "IMAGE_TOP_SUB", "size": 17, "bold": False, "line-height": 1},
        "[SUBJECT]": {"key": "SUBJECT", "size": 12, "bold": False}
     }},

    {"template": "./templates/blue-blur/dark/image-right.pptx",
     "image": "./templates/image-h.png",
     "blocks": {
        "[IMAGE_BOTTOM_RIGHT]": {"key": "IMAGE_BOTTOM_RIGHT", "bold": True, "size": 40, "line-height": 0.9},
        "[IMAGE_BOTTOM_RIGHT_SUB]": {"key": "IMAGE_BOTTOM_RIGHT_SUB", "size": 23, "bold": False, "line-height": 1},
        "[IMAGE_BOTTOM_RIGHT_CAP]": {"key": "IMAGE_BOTTOM_RIGHT_CAP", "size": 12},
        "[SUBJECT]": {"key": "SUBJECT", "size": 12, "bold": False}
     }},

    {"template": "./templates/blue-blur/dark/image-bottom.pptx",
    "image": "./templates/image-w2.png",
     "blocks": {
        "[IMAGE_BOTTOM]": {"key": "IMAGE_BOTTOM", "bold": True, "size": 40, "line-height": 0.9},
        "[IMAGE_BOTTOM_SUB]": {"key": "IMAGE_BOTTOM_SUB", "size": 23, "bold": False, "line-height": 1},
        "[SUBJECT]": {"key": "SUBJECT", "size": 12, "bold": False}
     }},

    {"template": "./templates/blue-blur/dark/Image-left.pptx",
    "image": "./templates/image-h2.png",
     "blocks": {
        "[IMAGE_BOTTOM_LEFT]": {"key": "IMAGE_BOTTOM_LEFT", "bold": True, "size": 40, "line-height": 0.9},
        "[IMAGE_BOTTOM_LEFT_SUB]": {"key": "IMAGE_BOTTOM_LEFT_SUB", "size": 17,  "bold": False, "line-height": 1},
        "[IMAGE_BOTTOM_LEFT_CAP]": {"key": "IMAGE_BOTTOM_LEFT_CAP", "size": 12},
         "[SUBJECT]": {"key": "SUBJECT", "size": 12, "bold": False}
     }},
    {"template": "./templates/blue-blur/dark/CTA.pptx",
      "image": None,
     "blocks": {
        "[CTA]": {"key": "CTA", "bold": True, "align": "center", "text-block-height": 200},
        "[CTA_SUB]": {"key": "CTA_SUB", "size": 17, "bold": False}
     }},

]

post_text = """
[SUBJECT]
Transforming Business with AI Agents
[HOOK]
Transforming Business with AI Agents

[HOOK_SUB]
Optimize resources, maximize ROI, and extend your B2B marketing budget through data-backed decisions.

[STORY]
Why this matters today

[STORY_SUB]
Early adopters of AI agents are already boosting efficiency, reducing costs, 
and creating new value streams. Those who delay risk being disrupted.

[TOPIC]
5 ways AI agents are reshaping business

[TOPIC_SUB]
1. Automating repetitive tasks  
2. Enabling contextual decision-making  
3. Personalizing customer engagement  
4. Unlocking new business models  
5. Driving cost savings at scale  

[IMAGE_TOP]
AI Agents in Action
[IMAGE_TOP_SUB]
Imagine a customer service agent that not only answers questions but anticipates needs and resolves issues proactively.

[IMAGE_BOTTOM_RIGHT]
From insight to execution

[IMAGE_BOTTOM_RIGHT_SUB]
Reasoning models analyze the “why” behind the data and recommend the best next step.

[IMAGE_BOTTOM_RIGHT_CAP]
⚡ Faster, smarter, more confident decisions

[IMAGE_BOTTOM]
Future-proof your company
[IMAGE_BOTTOM_SUB]
Reasoning models analyze the “why” behind the data and recommend the best next step.

[CTA]
What are your thoughts?
[CTA_SUB]
Ready to transform your  B2B marketing with data?

[IMAGE_BOTTOM_LEFT]
The competitive edge

[IMAGE_BOTTOM_LEFT_SUB]
Success will belong to those who combine data, reasoning, and execution.  

[IMAGE_BOTTOM_LEFT_CAP]
🚀 Reasoning is the new ROI
"""
build_carousel(post_text, template_mappings, "my_carousel.pptx")
merge_pptx_slides('./concluded', './concluded/done/production-ready_images.pptx')
