import os
import sys
from pptx import Presentation
from lxml import etree
import datetime
from pptx.parts.image import ImagePart
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
# Monkeypatch para dar nomes únicos com timestamp aos arquivos de mídia
def _custom_new_image_part(package, image, ext):
    ts = datetime.datetime.now().strftime("%Y%m%d_%H%M%S_%f")
    partname = f"/ppt/media/IMG_{ts}{ext.lower()}"
    return package._package.create_part(partname, image.content_type, image.blob)


def merge_pptx_slides(input_dir: str, output_file: str):
    """
    Merge multiple PPTX files into a single presentation.
    Preserves backgrounds, layouts, text, images, colors and formatting
    by copying slides exactly as in the originals.
    """
    ImagePart._new = staticmethod(_custom_new_image_part)

    files = sorted(f for f in os.listdir(input_dir) if f.endswith(".pptx"))
    if not files:
        raise FileNotFoundError(f"No PPTX files found in {input_dir}")

    # Create a new empty presentation from the first file (to preserve theme/masters)
    base_file = os.path.join(input_dir, files[0])
    merged_prs = Presentation(base_file)

    # Remove all slides from the base
    while len(merged_prs.slides) > 0:
        rId = merged_prs.slides._sldIdLst[0].rId
        merged_prs.part.drop_rel(rId)
        del merged_prs.slides._sldIdLst[0]

    # Process each file
    for filename in files:
        filepath = os.path.join(input_dir, filename)
        source_prs = Presentation(filepath)

        for source_slide in source_prs.slides:
            # Get the source slide layout
            source_layout = source_slide.slide_layout
            
            # Try to find matching layout in merged presentation
            try:
                layout_idx = source_prs.slide_layouts.index(source_layout)
                target_layout = merged_prs.slide_layouts[layout_idx]
            except (ValueError, IndexError):
                target_layout = merged_prs.slide_layouts[6]  # Blank layout fallback
            
            # Create new slide
            new_slide = merged_prs.slides.add_slide(target_layout)
            
            # Remove all default shapes from the new slide
            for shape in list(new_slide.shapes):
                sp = shape.element
                sp.getparent().remove(sp)
            
            # Map to track old rId -> new rId for images and other relationships
            rid_map = {}
            
            # Copy all relationships (images, charts, etc.) from source slide
            for rel in source_slide.part.rels.values():
                # Skip slide layout relationship (already handled)
                if rel.reltype == RT.SLIDE_LAYOUT:
                    continue
                
                try:
                    # Get the related part (image, chart, etc.)
                    related_part = rel.target_part
                    
                    # Add the related part to the new slide
                    if rel.reltype == RT.IMAGE:
                        # For images, copy the image data
                        new_rel = new_slide.part.relate_to(related_part, rel.reltype)
                        rid_map[rel.rId] = new_rel.rId
                    else:
                        # For other relationships, just create the relationship
                        try:
                            new_rel = new_slide.part.relate_to(related_part, rel.reltype)
                            rid_map[rel.rId] = new_rel.rId
                        except:
                            pass
                except Exception as e:
                    print(f"⚠️  Warning: Could not copy relationship {rel.rId}: {e}")
                    continue
            
            # Copy the entire slide structure from source
            for shape in source_slide.shapes:
                el = shape.element
                # Clone the element completely
                new_el = etree.fromstring(etree.tostring(el))
                
                # Update image references (rId) in the cloned element
                # Look for blip elements (images)
                nsmap = {
                    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
                    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'
                }
                
                # Find all image references and update their rIds
                for blip in new_el.findall('.//a:blip', namespaces=nsmap):
                    old_rid = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                    if old_rid and old_rid in rid_map:
                        blip.set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed', 
                                rid_map[old_rid])
                
                # Add the shape to the new slide
                new_slide.shapes._spTree.append(new_el)
            
            # Copy slide background
            try:
                source_bg = source_slide.element.cSld.bg
                if source_bg is not None:
                    # Clone background element
                    new_bg = etree.fromstring(etree.tostring(source_bg))
                    
                    # Update background image references if present
                    nsmap = {
                        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
                    }
                    for blip in new_bg.findall('.//a:blip', namespaces=nsmap):
                        old_rid = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                        if old_rid and old_rid in rid_map:
                            blip.set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed', 
                                    rid_map[old_rid])
                    
                    # Remove existing background if present
                    existing_bg = new_slide.element.cSld.find('.//p:bg', 
                        namespaces={'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'})
                    if existing_bg is not None:
                        existing_bg.getparent().remove(existing_bg)
                    # Add new background
                    new_slide.element.cSld.insert(0, new_bg)
            except AttributeError:
                pass  # No background to copy
            
            # Copy color map override if present
            try:
                source_clrmap = source_slide.element.find('.//p:clrMapOvr',
                    namespaces={'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'})
                if source_clrmap is not None:
                    new_clrmap = etree.fromstring(etree.tostring(source_clrmap))
                    existing_clrmap = new_slide.element.find('.//p:clrMapOvr',
                        namespaces={'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'})
                    if existing_clrmap is not None:
                        existing_clrmap.getparent().remove(existing_clrmap)
                    new_slide.element.append(new_clrmap)
            except:
                pass

        print(f"✅ Added {len(source_prs.slides)} slide(s) from {filename}")

    merged_prs.save(output_file)
    print(f"\n🎉 Final merged file saved as: {output_file}")
    print(f"📊 Total slides: {len(merged_prs.slides)}")


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python merge_templates.py ./templates/blue-blur/dark merged-dark.pptx")
        sys.exit(1)

    input_dir = sys.argv[1]
    output_file = sys.argv[2]
    
    if not os.path.exists(input_dir):
        print(f"❌ Error: Directory '{input_dir}' not found")
        sys.exit(1)
    
    merge_pptx_slides(input_dir, output_file)